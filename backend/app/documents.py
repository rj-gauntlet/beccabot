"""Document parsing for PDF, DOCX, TXT, PPTX, XLSX, and CSV files.

Extracts text and tables from all supported formats. No manual knowledge.json needed.
"""

import csv
import io
import os
from pathlib import Path

# Ensure Tesseract is on PATH for OCR (common Windows install location)
_TESSERACT_PATH = Path(r"C:\Program Files\Tesseract-OCR")
if _TESSERACT_PATH.exists():
    os.environ["PATH"] = str(_TESSERACT_PATH) + os.pathsep + os.environ.get("PATH", "")

from docx import Document as DocxDocument
from pptx import Presentation

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".doc", ".txt", ".pptx", ".ppt", ".xlsx", ".csv"}


def parse_document(file_path: Path) -> str:
    """Extract text and tables from a supported document. Returns concatenated text."""
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return _parse_pdf(path)
    elif suffix in (".docx", ".doc"):
        return _parse_docx(path)
    elif suffix == ".txt":
        return _parse_txt(path)
    elif suffix in (".pptx", ".ppt"):
        return _parse_pptx(path)
    elif suffix == ".xlsx":
        return _parse_excel(path)
    elif suffix == ".csv":
        return _parse_csv(path)
    else:
        raise ValueError(f"Unsupported file type: {suffix}")


def _text_looks_like_garbage(text: str) -> bool:
    """Detect gibberish from corrupted PDF text layers (single chars, symbols, OCR noise)."""
    if not text or len(text.strip()) < 50:
        return True
    t = text.strip()
    # Low alphabetic ratio = mostly symbols/garbage
    alpha = sum(1 for c in t if c.isalpha())
    if alpha / len(t) < 0.5:
        return True
    # Need several substantive words (5+ chars) - filters OCR noise like "oss", "uae"
    words = [w for w in t.split() if len(w) >= 5 and sum(c.isalpha() for c in w) >= 4]
    if len(words) < 5:
        return True
    return False


def _parse_pdf(path: Path) -> str:
    """Extract text from PDF. Uses PyMuPDF, pypdf fallback, then OCR for scanned/image PDFs."""
    text = _parse_pdf_standard(path)
    if text and len(text.strip()) >= 100 and not _text_looks_like_garbage(text):
        return text

    # Try OCR for image-based/scanned PDFs (requires Tesseract installed)
    ocr_text = _parse_pdf_ocr(path)
    if ocr_text and not _text_looks_like_garbage(ocr_text):
        return ocr_text
    # Don't return garbage from corrupted text layers; return OCR or empty
    if ocr_text:
        return ocr_text
    if text and _text_looks_like_garbage(text):
        return ""  # Standard text is garbage, OCR failed; avoid indexing junk
    return text or ""


def _parse_pdf_standard(path: Path) -> str:
    """Extract text and tables using PyMuPDF, then pypdf fallback."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(path)
        parts = []
        for page in doc:
            # Standard text extraction
            t = page.get_text()
            if t:
                parts.append(t.strip())
            # Table extraction (addresses, headers, tabular data often missed by get_text)
            try:
                tbf = page.find_tables()
                if tbf:
                    for tbl in tbf.tables:
                        try:
                            rows = tbl.extract()
                            if rows:
                                table_lines = [
                                    " | ".join(str(c or "").strip() for c in row)
                                    for row in rows
                                ]
                                if table_lines:
                                    parts.append("\n".join(table_lines))
                        except Exception:
                            pass
            except Exception:
                pass
        doc.close()
        if parts:
            return "\n\n".join(parts)
    except Exception:
        pass
    try:
        from pypdf import PdfReader
        reader = PdfReader(path)
        parts = [p.extract_text() for p in reader.pages if p.extract_text()]
        return "\n\n".join(parts) if parts else ""
    except Exception:
        return ""


def _parse_pdf_ocr(path: Path) -> str:
    """OCR image-based PDFs using PyMuPDF + Tesseract. Returns empty if Tesseract not installed."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(path)
        parts = []
        for page in doc:
            # full=True: render page to pixmap and OCR it, bypassing corrupted text layers
            tp = page.get_textpage_ocr(full=True, dpi=150)
            text = page.get_text(textpage=tp) if tp else ""
            if text and text.strip():
                parts.append(text.strip())
        doc.close()
        return "\n\n".join(parts) if parts else ""
    except Exception:
        return ""


def _parse_docx(path: Path) -> str:
    doc = DocxDocument(path)
    parts = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text.strip())
    # Extract tables (addresses, schedules, etc.)
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [str(cell.text or "").strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            parts.append("\n".join(rows))
    return "\n\n".join(parts)


def _parse_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _parse_pptx(path: Path) -> str:
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
            if getattr(shape, "has_table", False) and shape.has_table:
                rows = []
                for row in shape.table.rows:
                    cells = [str(cell.text or "").strip() for cell in row.cells]
                    rows.append(" | ".join(cells))
                if rows:
                    slide_text.append("\n".join(rows))
        if slide_text:
            parts.append("\n".join(slide_text))
    return "\n\n".join(parts)


def _parse_excel(path: Path) -> str:
    """Extract text from Excel (.xlsx). Uses openpyxl."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        parts = []
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c or "").strip() for c in row]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                parts.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(parts)
    except ImportError:
        raise ValueError(
            "Excel support requires openpyxl. Install with: pip install openpyxl"
        )


def _parse_csv(path: Path) -> str:
    """Extract text from CSV."""
    text = path.read_text(encoding="utf-8", errors="replace")
    parts = []
    for row in csv.reader(io.StringIO(text)):
        if row:
            parts.append(" | ".join(str(c or "").strip() for c in row))
    return "\n".join(parts)
